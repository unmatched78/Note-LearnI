# core/utils/file_processing.py
import logging
from PyPDF2 import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation
from langchain.text_splitter import RecursiveCharacterTextSplitter

logger = logging.getLogger(__name__)

def extract_text_from_file(uploaded_file):
    """
    Extract text from an uploaded file based on its type.
    Supports PDF, DOCX, TXT, and PPTX files.
    Returns a tuple of (text, file_name).
    """
    file_name = uploaded_file.name
    logger.debug("Extracting text from file: %s", file_name)
    file_type = file_name.split('.')[-1].lower()

    try:
        if file_type == "pdf":
            reader = PdfReader(uploaded_file)
            text = "\n".join(page.extract_text() for page in reader.pages if page.extract_text())
            return text, file_name
        elif file_type == "docx":
            doc = DocxDocument(uploaded_file)
            text = "\n".join(paragraph.text for paragraph in doc.paragraphs)
            return text, file_name
        elif file_type == "txt":
            text = uploaded_file.read().decode("utf-8")
            return text, file_name
        elif file_type == "pptx":
            prs = Presentation(uploaded_file)
            text_list = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_list.append(shape.text)
            text = "\n".join(text_list)
            return text, file_name
        else:
            logger.error("Unsupported file type: %s", file_type)
            return None, file_name
    except Exception as e:
        logger.error("Error extracting text from file %s: %s", file_name, e)
        return None, file_name

def split_text_into_chunks(text, chunk_size=1000, chunk_overlap=100):
    """
    Split text into smaller chunks using RecursiveCharacterTextSplitter.
    """
    if not text:
        return []
    splitter = RecursiveCharacterTextSplitter(chunk_size=chunk_size, chunk_overlap=chunk_overlap)
    chunks = splitter.split_text(text)
    return chunks